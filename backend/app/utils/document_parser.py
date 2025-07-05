import os
import logging
from typing import Dict, Any, Optional
from pathlib import Path

# Document processing libraries
try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from app.models.document import FileType

logger = logging.getLogger(__name__)


class DocumentParser:
    """Parser for various document formats."""
    
    def __init__(self):
        self.supported_formats = {
            FileType.PDF: self._parse_pdf,
            FileType.DOCX: self._parse_docx,
            FileType.PPTX: self._parse_pptx,
            FileType.TXT: self._parse_txt,
            FileType.MD: self._parse_md
        }
    
    async def parse_document(self, file_path: str, file_type: FileType) -> Dict[str, Any]:
        """
        Parse a document and extract text and metadata.
        
        Args:
            file_path: Path to the document file
            file_type: Type of the document
            
        Returns:
            Dictionary with extracted text and metadata
        """
        try:
            if file_type not in self.supported_formats:
                raise ValueError(f"Unsupported file type: {file_type}")
            
            parser_func = self.supported_formats[file_type]
            result = await parser_func(file_path)
            
            # Add common metadata
            result["metadata"]["file_type"] = file_type.value
            result["metadata"]["file_size"] = os.path.getsize(file_path)
            result["metadata"]["file_name"] = os.path.basename(file_path)
            
            return result
            
        except Exception as e:
            logger.error(f"Error parsing document {file_path}: {str(e)}")
            raise
    
    async def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """Parse PDF document."""
        if not PDF_AVAILABLE:
            raise RuntimeError("PyPDF2 not installed. Install with: pip install PyPDF2")
        
        try:
            text_content = []
            metadata = {}
            
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                
                # Extract metadata
                if pdf_reader.metadata:
                    metadata.update({
                        "title": pdf_reader.metadata.get("/Title", ""),
                        "author": pdf_reader.metadata.get("/Author", ""),
                        "subject": pdf_reader.metadata.get("/Subject", ""),
                        "creator": pdf_reader.metadata.get("/Creator", ""),
                        "producer": pdf_reader.metadata.get("/Producer", ""),
                        "creation_date": str(pdf_reader.metadata.get("/CreationDate", "")),
                        "modification_date": str(pdf_reader.metadata.get("/ModDate", ""))
                    })
                
                metadata["pages"] = len(pdf_reader.pages)
                
                # Extract text from all pages
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(f"[Page {page_num + 1}]\n{page_text}")
            
            return {
                "text": "\n\n".join(text_content),
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Error parsing PDF: {str(e)}")
            raise
    
    async def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """Parse DOCX document."""
        if not DOCX_AVAILABLE:
            raise RuntimeError("python-docx not installed. Install with: pip install python-docx")
        
        try:
            doc = DocxDocument(file_path)
            
            # Extract text
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract metadata
            metadata = {}
            if doc.core_properties:
                metadata.update({
                    "title": doc.core_properties.title or "",
                    "author": doc.core_properties.author or "",
                    "subject": doc.core_properties.subject or "",
                    "keywords": doc.core_properties.keywords or "",
                    "comments": doc.core_properties.comments or "",
                    "created": str(doc.core_properties.created) if doc.core_properties.created else "",
                    "modified": str(doc.core_properties.modified) if doc.core_properties.modified else "",
                    "last_modified_by": doc.core_properties.last_modified_by or ""
                })
            
            metadata["paragraphs"] = len(doc.paragraphs)
            metadata["tables"] = len(doc.tables)
            
            # Extract table content
            table_content = []
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    table_text.append(" | ".join(row_text))
                
                if table_text:
                    table_content.append("\n".join(table_text))
            
            # Combine text and table content
            full_text = "\n\n".join(text_content)
            if table_content:
                full_text += "\n\n[TABLES]\n" + "\n\n".join(table_content)
            
            return {
                "text": full_text,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Error parsing DOCX: {str(e)}")
            raise
    
    async def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """Parse PPTX document."""
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not installed. Install with: pip install python-pptx")
        
        try:
            presentation = Presentation(file_path)
            
            text_content = []
            metadata = {}
            
            # Extract metadata
            if presentation.core_properties:
                metadata.update({
                    "title": presentation.core_properties.title or "",
                    "author": presentation.core_properties.author or "",
                    "subject": presentation.core_properties.subject or "",
                    "keywords": presentation.core_properties.keywords or "",
                    "comments": presentation.core_properties.comments or "",
                    "created": str(presentation.core_properties.created) if presentation.core_properties.created else "",
                    "modified": str(presentation.core_properties.modified) if presentation.core_properties.modified else "",
                    "last_modified_by": presentation.core_properties.last_modified_by or ""
                })
            
            metadata["slides"] = len(presentation.slides)
            
            # Extract text from all slides
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = []
                slide_text.append(f"[Slide {slide_num + 1}]")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.append("\n".join(slide_text))
            
            return {
                "text": "\n\n".join(text_content),
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Error parsing PPTX: {str(e)}")
            raise
    
    async def _parse_txt(self, file_path: str) -> Dict[str, Any]:
        """Parse plain text document."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            # Basic metadata
            metadata = {
                "lines": len(content.splitlines()),
                "characters": len(content),
                "words": len(content.split())
            }
            
            return {
                "text": content,
                "metadata": metadata
            }
            
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    content = file.read()
                
                metadata = {
                    "lines": len(content.splitlines()),
                    "characters": len(content),
                    "words": len(content.split()),
                    "encoding": "latin-1"
                }
                
                return {
                    "text": content,
                    "metadata": metadata
                }
            except Exception as e:
                logger.error(f"Error parsing TXT with latin-1 encoding: {str(e)}")
                raise
        except Exception as e:
            logger.error(f"Error parsing TXT: {str(e)}")
            raise
    
    async def _parse_md(self, file_path: str) -> Dict[str, Any]:
        """Parse Markdown document."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            # Basic metadata
            metadata = {
                "lines": len(content.splitlines()),
                "characters": len(content),
                "words": len(content.split())
            }
            
            # Count markdown elements
            lines = content.splitlines()
            headings = sum(1 for line in lines if line.strip().startswith('#'))
            code_blocks = content.count('```')
            links = content.count('[') + content.count('](')
            
            metadata.update({
                "headings": headings,
                "code_blocks": code_blocks // 2,  # Divide by 2 for opening/closing pairs
                "links": links
            })
            
            return {
                "text": content,
                "metadata": metadata
            }
            
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    content = file.read()
                
                metadata = {
                    "lines": len(content.splitlines()),
                    "characters": len(content),
                    "words": len(content.split()),
                    "encoding": "latin-1"
                }
                
                return {
                    "text": content,
                    "metadata": metadata
                }
            except Exception as e:
                logger.error(f"Error parsing MD with latin-1 encoding: {str(e)}")
                raise
        except Exception as e:
            logger.error(f"Error parsing MD: {str(e)}")
            raise
    
    def get_supported_formats(self) -> list:
        """Get list of supported file formats."""
        return list(self.supported_formats.keys())
    
    def is_format_supported(self, file_type: FileType) -> bool:
        """Check if a file format is supported."""
        return file_type in self.supported_formats 